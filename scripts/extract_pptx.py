"""Extract text from PPTX file."""
import sys
from pptx import Presentation

filepath = sys.argv[1]
prs = Presentation(filepath)
texts = []
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
print("\n".join(texts))
